"""
build_pptx.py
Generates return-verification.pptx from scratch using python-pptx.
Run from the project directory: python3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy
import os

# ---------------------------------------------------------------------------
# CONSTANTS
# ---------------------------------------------------------------------------

W = Inches(13.333)   # 16:9 width
H = Inches(7.5)      # 16:9 height

# Brand colors
TEAL        = RGBColor(0,   140, 149)
TEAL_DARK   = RGBColor(0,   107, 114)
TEAL_LIGHT  = RGBColor(221, 240, 238)   # foam
GOLDEN      = RGBColor(214, 154,  45)
GOLDEN_LIGHT= RGBColor(255, 237, 186)
SLATE       = RGBColor( 76,  76,  78)
SLATE_LIGHT = RGBColor(150, 150, 152)
FOAM        = RGBColor(221, 240, 238)
FOG         = RGBColor(237, 236, 232)
WHITE       = RGBColor(255, 255, 255)
RED_SOFT    = RGBColor(192,  57,  43)

LOGO_PATH   = "../USEFULL-icons/USEFULL-Icon-Registered_Color.svg"
ICON_PNG    = "../USEFULL-icons/icon-512.png"   # fallback raster icon for pptx
STATION_IMG = "../return-station-concept.png"

# Typography
FONT = "Calibri"          # available on most systems
FONT_DISPLAY = "Georgia"  # serif for headings


# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    """Return the blank slide layout."""
    return prs.slide_layouts[6]


def add_slide(prs):
    return prs.slides.add_slide(blank_layout(prs))


def fill_bg(slide, color):
    """Solid background fill on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_name=FONT, font_size=Pt(14),
                bold=False, italic=False, color=SLATE, align=PP_ALIGN.LEFT,
                wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing.pt * 100)))
    return txBox


def add_label(slide, x, y, w, text, color=GOLDEN):
    """Small all-caps label / eyebrow."""
    return add_textbox(slide, x, y, w, Inches(0.3), text.upper(),
                       font_size=Pt(10), bold=True, color=color)


def add_heading(slide, x, y, w, text, font_size=Pt(32), color=TEAL):
    return add_textbox(slide, x, y, w, Inches(0.8), text,
                       font_name=FONT_DISPLAY, font_size=font_size,
                       color=color)


def add_body(slide, x, y, w, h, text, font_size=Pt(14), color=SLATE,
             bold=False, italic=False, align=PP_ALIGN.LEFT):
    return add_textbox(slide, x, y, w, h, text,
                       font_size=font_size, color=color,
                       bold=bold, italic=italic, align=align)


def add_multiline_textbox(slide, x, y, w, h, paragraphs, default_font=FONT,
                          default_size=Pt(14), default_color=SLATE):
    """
    paragraphs: list of dicts with keys:
      text, font_name, font_size, bold, italic, color, align, space_before
    """
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for para in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = para.get('align', PP_ALIGN.LEFT)
        if 'space_before' in para:
            p.space_before = para['space_before']
        run = p.add_run()
        run.text = para.get('text', '')
        run.font.name = para.get('font_name', default_font)
        run.font.size = para.get('font_size', default_size)
        run.font.bold = para.get('bold', False)
        run.font.italic = para.get('italic', False)
        run.font.color.rgb = para.get('color', default_color)
    return txBox


def bullet_block(slide, x, y, w, items, font_size=Pt(13), color=SLATE,
                 bullet_color=TEAL, indent=Inches(0.2), row_height=Inches(0.35)):
    """Draw a simple bulleted list as individual shapes."""
    cy = y
    for item in items:
        # bullet dot
        dot_size = Inches(0.07)
        dot = add_rect(slide,
                       x, cy + Inches(0.12),
                       dot_size, dot_size,
                       fill_color=bullet_color)
        dot.line.fill.background()
        # text
        add_textbox(slide, x + indent, cy, w - indent, row_height,
                    item, font_size=font_size, color=color)
        cy += row_height
    return cy   # returns bottom y


def add_divider_line(slide, x, y, w, color=GOLDEN, thickness=Pt(2)):
    line = slide.shapes.add_shape(1, x, y, w, Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_logo(slide, x=None, y=None, h=Inches(0.3)):
    """Add the icon PNG as a small logo. Positioned top-right if x/y omitted."""
    if not os.path.exists(ICON_PNG):
        return
    if x is None:
        x = W - Inches(0.9)
    if y is None:
        y = Inches(0.15)
    slide.shapes.add_picture(ICON_PNG, x, y, height=h)


def slide_num_label(slide, num, total=11, color=SLATE_LIGHT):
    add_textbox(slide,
                W - Inches(1.4), Inches(0.18), Inches(1.2), Inches(0.3),
                f"{num:02d} / {total:02d}",
                font_size=Pt(9), color=color, align=PP_ALIGN.RIGHT)


def top_bar(slide, num, bg_is_dark=False):
    """Logo icon + slide number in top bar."""
    add_logo(slide, x=Inches(0.45), y=Inches(0.18), h=Inches(0.28))
    slide_num_label(slide, num, color=SLATE_LIGHT if not bg_is_dark else RGBColor(200,200,200))


def add_teal_accent_bar(slide, x, y, h, w=Inches(0.05)):
    """Vertical teal accent bar (left callout border)."""
    bar = add_rect(slide, x, y, w, h, fill_color=TEAL)
    bar.line.fill.background()
    return bar


def add_golden_accent_bar(slide, x, y, h, w=Inches(0.05)):
    bar = add_rect(slide, x, y, w, h, fill_color=GOLDEN)
    bar.line.fill.background()
    return bar


# ---------------------------------------------------------------------------
# CONTENT MARGIN
# ---------------------------------------------------------------------------
ML = Inches(0.55)    # left margin
MR = Inches(0.55)    # right margin
MT = Inches(0.7)     # top (below header)
CONTENT_W = W - ML - MR
CONTENT_TOP = Inches(1.1)   # below top bar


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    s = add_slide(prs)
    fill_bg(s, FOAM)

    # Decorative teal block top-right
    deco_w = Inches(4.5)
    deco_h = Inches(4.5)
    deco = add_rect(s, W - deco_w, Inches(0), deco_w, deco_h,
                    fill_color=RGBColor(0, 140, 149))
    deco.fill.fore_color.rgb = TEAL
    # make it semi-transparent by using a lighter tint
    deco2 = add_rect(s, W - Inches(3), Inches(0), Inches(3), Inches(3),
                     fill_color=RGBColor(180, 225, 227))
    deco2.line.fill.background()

    top_bar(s, 1)

    # Eyebrow
    add_label(s, ML, Inches(2.2), Inches(6), "Feasibility Research")

    # Title
    title_box = add_textbox(
        s, ML, Inches(2.55), Inches(7.5), Inches(1.8),
        "Return Verification System",
        font_name=FONT_DISPLAY, font_size=Pt(44), bold=False, color=TEAL
    )

    # Golden divider
    add_divider_line(s, ML, Inches(4.5), Inches(3.2))

    # Subtitle
    add_body(s, ML, Inches(4.7), Inches(7), Inches(0.5),
             "One potential build plan — AI camera verification for USEFULL return bins",
             font_size=Pt(14), color=SLATE_LIGHT, italic=True)

    # Meta
    add_body(s, ML, Inches(5.2), Inches(5), Inches(0.4),
             "Owen Barron  ·  March 2026",
             font_size=Pt(13), color=SLATE_LIGHT)

    return s


def slide_02_problem(prs):
    s = add_slide(prs)
    fill_bg(s, WHITE)
    top_bar(s, 2)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Starting Point")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "The Problem We Were Asked to Solve", font_size=Pt(30))

    add_body(s, ML, CONTENT_TOP + Inches(1.05), CONTENT_W, Inches(0.45),
             "Alison's ask: can we use an AI camera to verify containers are actually returned after a QR scan?",
             font_size=Pt(14))

    # Flow diagram — 3 boxes connected by arrows
    flow_y = CONTENT_TOP + Inches(1.65)
    flow_h = Inches(1.1)
    box_w  = Inches(3.2)
    gap_w  = Inches(3.8)   # "gap" box is slightly wider
    arrow_w = Inches(0.5)
    total_flow_w = box_w + arrow_w + box_w + arrow_w + gap_w
    flow_x = ML

    # Box 1
    b1 = add_rect(s, flow_x, flow_y, box_w, flow_h, fill_color=TEAL)
    b1.line.fill.background()
    add_multiline_textbox(s, flow_x + Inches(0.15), flow_y + Inches(0.12),
                          box_w - Inches(0.3), flow_h - Inches(0.2),
                          [{'text': '1', 'font_name': FONT_DISPLAY, 'font_size': Pt(28),
                            'color': RGBColor(180,225,227), 'bold': False},
                           {'text': 'Student scans QR code', 'font_size': Pt(13),
                            'color': WHITE, 'space_before': Pt(2)}])

    # Arrow 1
    add_textbox(s, flow_x + box_w, flow_y, arrow_w, flow_h,
                "→", font_size=Pt(22), color=TEAL, align=PP_ALIGN.CENTER)

    # Box 2
    b2x = flow_x + box_w + arrow_w
    b2 = add_rect(s, b2x, flow_y, box_w, flow_h, fill_color=TEAL)
    b2.line.fill.background()
    add_multiline_textbox(s, b2x + Inches(0.15), flow_y + Inches(0.12),
                          box_w - Inches(0.3), flow_h - Inches(0.2),
                          [{'text': '2', 'font_name': FONT_DISPLAY, 'font_size': Pt(28),
                            'color': RGBColor(180,225,227), 'bold': False},
                           {'text': 'System expects a return', 'font_size': Pt(13),
                            'color': WHITE, 'space_before': Pt(2)}])

    # Arrow 2
    arrow2x = b2x + box_w
    add_textbox(s, arrow2x, flow_y, arrow_w, flow_h,
                "→", font_size=Pt(22), color=RED_SOFT, align=PP_ALIGN.CENTER)

    # Gap box
    gapx = arrow2x + arrow_w
    gap = add_rect(s, gapx, flow_y, gap_w, flow_h,
                   fill_color=RGBColor(255, 240, 238),
                   line_color=RED_SOFT, line_width=Pt(1.5))
    add_multiline_textbox(s, gapx + Inches(0.15), flow_y + Inches(0.08),
                          gap_w - Inches(0.3), flow_h - Inches(0.15),
                          [{'text': '?  Gap', 'font_name': FONT_DISPLAY, 'font_size': Pt(22),
                            'color': RED_SOFT, 'bold': False},
                           {'text': 'Did a real container actually drop?',
                            'font_size': Pt(13), 'color': RED_SOFT,
                            'space_before': Pt(3)}])

    # Callout
    callout_y = flow_y + flow_h + Inches(0.4)
    add_golden_accent_bar(s, ML, callout_y, Inches(0.55))
    add_body(s, ML + Inches(0.15), callout_y, CONTENT_W - Inches(0.15), Inches(0.6),
             "Current vulnerability: a student can scan the QR code, walk away without dropping, "
             "or drop trash to trigger a simple sensor — and receive credit.",
             font_size=Pt(13), color=SLATE_LIGHT, italic=True)

    return s


def slide_03_right_problem(prs):
    s = add_slide(prs)
    fill_bg(s, FOG)
    top_bar(s, 3)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Field Evidence")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "But Is This the Right Problem?", font_size=Pt(30))

    col_y = CONTENT_TOP + Inches(1.05)
    col_h = Inches(4.5)
    col_w = (CONTENT_W - Inches(0.3)) / 2

    # Left column — solves
    lx = ML
    solve_bg = add_rect(s, lx, col_y, col_w, col_h,
                        fill_color=RGBColor(220, 242, 240),
                        line_color=TEAL, line_width=Pt(1))
    add_textbox(s, lx + Inches(0.2), col_y + Inches(0.15), col_w - Inches(0.4), Inches(0.35),
                "WHAT CAMERA VERIFICATION SOLVES",
                font_size=Pt(9), bold=True, color=TEAL)
    bullet_block(s, lx + Inches(0.2), col_y + Inches(0.55), col_w - Inches(0.4),
                 ["Scan-and-walk-away (QR scanned, nothing dropped)",
                  "Fake drops (trash triggered instead of container)",
                  "Verifiable return event in the system of record"],
                 font_size=Pt(13), bullet_color=TEAL)

    # Right column — doesn't solve
    rx = ML + col_w + Inches(0.3)
    add_rect(s, rx, col_y, col_w, col_h,
             fill_color=RGBColor(253, 246, 225),
             line_color=GOLDEN, line_width=Pt(1))
    add_textbox(s, rx + Inches(0.2), col_y + Inches(0.15), col_w - Inches(0.4), Inches(0.35),
                "WHAT IT DOESN'T SOLVE",
                font_size=Pt(9), bold=True, color=GOLDEN)
    bullet_block(s, rx + Inches(0.2), col_y + Inches(0.55), col_w - Inches(0.4),
                 ["Containers never brought back at all (NAU's reported problem)",
                  "Student-employee collusion to skip checkout entirely",
                  "Student protest and non-participation (UMass Lowell)",
                  "Loss at other points in the container lifecycle"],
                 font_size=Pt(13), bullet_color=GOLDEN)

    # Callout bottom
    callout_y = col_y + col_h + Inches(0.2)
    add_golden_accent_bar(s, ML, callout_y, Inches(0.5))
    add_body(s, ML + Inches(0.15), callout_y, CONTENT_W, Inches(0.5),
             "A camera system solves a specific failure mode. Before investing, we should confirm "
             "this failure mode is actually driving losses.",
             font_size=Pt(13), color=SLATE, italic=True)

    return s


def slide_04_joshua(prs):
    s = add_slide(prs)
    fill_bg(s, WHITE)
    top_bar(s, 4)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Expert Guidance")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "What Joshua Told Us", font_size=Pt(30))

    add_body(s, ML, CONTENT_TOP + Inches(1.0), CONTENT_W, Inches(0.45),
             "Joshua is CEO of Amherst Intelligence Security, referred by Igor. After advising on the concept, "
             "he offered to work with us directly in his personal capacity — he's genuinely excited by the problem.",
             font_size=Pt(12), color=SLATE_LIGHT, italic=True)

    # 4 insight cards: 2×2 grid
    card_w = (CONTENT_W - Inches(0.3)) / 2
    card_h = Inches(1.9)
    card_gap = Inches(0.15)
    grid_y = CONTENT_TOP + Inches(1.6)

    cards = [
        ("01", "Don't point into the bin",
         "A messy pile of wet containers is a dead end. Focus on the throat — where the scene can be controlled."),
        ("02", "More cameras beats better cameras",
         "Multiple inexpensive cameras from more angles outperforms a single high-end unit."),
        ("03", "You need a lot of data",
         "This is a scaling problem. The model needs thousands of real return events before it's useful."),
        ("04", "Expect a long iteration cycle",
         "Getting something that works even a little in a real food-court environment takes significant time."),
    ]

    positions = [
        (ML, grid_y),
        (ML + card_w + card_gap, grid_y),
        (ML, grid_y + card_h + card_gap),
        (ML + card_w + card_gap, grid_y + card_h + card_gap),
    ]

    for (num, title, body), (cx, cy) in zip(cards, positions):
        card = add_rect(s, cx, cy, card_w, card_h,
                        fill_color=RGBColor(247, 252, 252),
                        line_color=RGBColor(200, 230, 232), line_width=Pt(1))
        # top border teal
        top_line = add_rect(s, cx, cy, card_w, Inches(0.06), fill_color=TEAL)
        top_line.line.fill.background()
        add_textbox(s, cx + Inches(0.15), cy + Inches(0.1), Inches(0.5), Inches(0.5),
                    num, font_name=FONT_DISPLAY, font_size=Pt(24), color=RGBColor(180,225,227))
        add_textbox(s, cx + Inches(0.15), cy + Inches(0.55), card_w - Inches(0.3), Inches(0.4),
                    title, font_size=Pt(13), bold=True, color=SLATE)
        add_textbox(s, cx + Inches(0.15), cy + Inches(0.95), card_w - Inches(0.3), Inches(0.9),
                    body, font_size=Pt(12), color=SLATE_LIGHT)

    # Design principle callout
    dp_y = grid_y + card_h * 2 + card_gap * 2 + Inches(0.1)
    add_teal_accent_bar(s, ML, dp_y, Inches(0.35))
    add_body(s, ML + Inches(0.15), dp_y, CONTENT_W, Inches(0.35),
             "Design principle: Watch the passage, not the pile. The throat concept on the next slide was built on this advice.",
             font_size=Pt(12), bold=True, color=TEAL_DARK)

    return s


def slide_05_throat(prs):
    s = add_slide(prs)
    fill_bg(s, FOAM)
    top_bar(s, 5)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Sensing Architecture — Developed with Joshua")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "The Throat Concept", font_size=Pt(30))

    # Left: diagram (described as labeled shapes)
    diag_x = ML
    diag_y = CONTENT_TOP + Inches(1.05)
    diag_w = Inches(6.5)
    diag_h = Inches(5.2)

    # Diagram background
    add_rect(s, diag_x, diag_y, diag_w, diag_h,
             fill_color=RGBColor(247, 252, 251),
             line_color=RGBColor(200, 230, 232), line_width=Pt(1))

    # Outer bin walls (dashed — simulate with rect + no fill)
    bin_margin = Inches(0.4)
    add_rect(s, diag_x + bin_margin, diag_y + Inches(0.4),
             diag_w - bin_margin*2, diag_h - Inches(0.8),
             line_color=RGBColor(154, 173, 164), line_width=Pt(1))

    # Left rail
    rail_x = diag_x + bin_margin + Inches(0.1)
    rail_y = diag_y + Inches(1.0)
    rail_w = Inches(0.7)
    rail_h = Inches(3.0)
    add_rect(s, rail_x, rail_y, rail_w, rail_h,
             fill_color=RGBColor(220, 242, 240),
             line_color=TEAL, line_width=Pt(1.5))
    add_textbox(s, rail_x, rail_y + Inches(1.1), rail_w, Inches(0.8),
                "LEFT\nRAIL\n1.75\"", font_size=Pt(8), bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    # Right rail
    rrail_x = diag_x + diag_w - bin_margin - Inches(0.1) - rail_w
    add_rect(s, rrail_x, rail_y, rail_w, rail_h,
             fill_color=RGBColor(220, 242, 240),
             line_color=TEAL, line_width=Pt(1.5))
    add_textbox(s, rrail_x, rail_y + Inches(1.1), rail_w, Inches(0.8),
                "RIGHT\nRAIL\n1.75\"", font_size=Pt(8), bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    # Throat clear zone
    throat_x = rail_x + rail_w
    throat_w = rrail_x - throat_x
    add_rect(s, throat_x, rail_y, throat_w, rail_h,
             fill_color=RGBColor(235, 248, 246))
    add_textbox(s, throat_x, rail_y + Inches(1.2), throat_w, Inches(0.6),
                "CLEAR THROAT\n10.5\"", font_size=Pt(9), bold=True,
                color=SLATE, align=PP_ALIGN.CENTER)

    # Zone label above throat
    add_textbox(s, throat_x, rail_y - Inches(0.35), throat_w, Inches(0.35),
                "ENTRY → PASSAGE → COMMIT",
                font_size=Pt(8), color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

    # Camera dots
    cam_r = Inches(0.18)
    # A - left rail, upper
    cam_ax = rail_x + rail_w/2 - cam_r
    cam_ay = rail_y + Inches(0.4)
    add_rect(s, cam_ax, cam_ay, cam_r*2, cam_r*2,
             fill_color=TEAL, line_color=WHITE, line_width=Pt(1))
    add_textbox(s, cam_ax - Inches(0.35), cam_ay - Inches(0.02), Inches(0.3), Inches(0.25),
                "A", font_size=Pt(10), bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

    # B - right rail, upper
    cam_bx = rrail_x + rail_w/2 - cam_r
    add_rect(s, cam_bx, cam_ay, cam_r*2, cam_r*2,
             fill_color=TEAL, line_color=WHITE, line_width=Pt(1))
    add_textbox(s, cam_bx + cam_r*2 + Inches(0.05), cam_ay - Inches(0.02),
                Inches(0.3), Inches(0.25), "B", font_size=Pt(10), bold=True, color=TEAL)

    # C - left rail, lower
    cam_cy_pos = rail_y + Inches(2.1)
    add_rect(s, cam_ax, cam_cy_pos, cam_r*2, cam_r*2,
             fill_color=GOLDEN, line_color=WHITE, line_width=Pt(1))
    add_textbox(s, cam_ax - Inches(0.35), cam_cy_pos - Inches(0.02),
                Inches(0.3), Inches(0.25), "C", font_size=Pt(10), bold=True,
                color=GOLDEN, align=PP_ALIGN.RIGHT)

    # D - right rail, lower
    add_rect(s, cam_bx, cam_cy_pos, cam_r*2, cam_r*2,
             fill_color=GOLDEN, line_color=WHITE, line_width=Pt(1))
    add_textbox(s, cam_bx + cam_r*2 + Inches(0.05), cam_cy_pos - Inches(0.02),
                Inches(0.3), Inches(0.25), "D", font_size=Pt(10), bold=True, color=GOLDEN)

    # 14" brace label
    add_textbox(s, diag_x + bin_margin, diag_y + diag_h - Inches(0.45),
                diag_w - bin_margin*2, Inches(0.35),
                "←————— 14\" slot envelope —————→",
                font_size=Pt(9), color=SLATE, align=PP_ALIGN.CENTER)

    # Right: legend
    leg_x = diag_x + diag_w + Inches(0.3)
    leg_w = W - leg_x - MR
    leg_y = diag_y

    add_textbox(s, leg_x, leg_y, leg_w, Inches(0.35),
                "4 cameras, 2 rails", font_size=Pt(13), bold=True, color=SLATE)

    legend_items = [
        (TEAL,   "A, B — Cross-view main cameras"),
        (GOLDEN, "C, D — Entry / commit watchers"),
    ]
    ly = leg_y + Inches(0.45)
    for lcolor, ltext in legend_items:
        dot = add_rect(s, leg_x, ly + Inches(0.07), Inches(0.14), Inches(0.14),
                       fill_color=lcolor)
        dot.line.fill.background()
        add_textbox(s, leg_x + Inches(0.22), ly, leg_w - Inches(0.22), Inches(0.3),
                    ltext, font_size=Pt(12), color=SLATE)
        ly += Inches(0.38)

    add_textbox(s, leg_x, ly + Inches(0.1), leg_w, Inches(0.35),
                "Cameras sealed behind polycarbonate windows — no electronics in the wet zone",
                font_size=Pt(11), color=SLATE_LIGHT, italic=True)

    ly += Inches(0.6)
    add_textbox(s, leg_x, ly, leg_w, Inches(0.35),
                "Detection question:", font_size=Pt(12), bold=True, color=SLATE)
    add_body(s, leg_x, ly + Inches(0.35), leg_w, Inches(0.7),
             "Did something substantial pass fully through the throat after a QR scan?",
             font_size=Pt(12), color=SLATE, italic=True)

    ly += Inches(1.2)
    bullets = ["USB cameras only — no bulky electronics in wet zone",
               "Scan-triggered capture window",
               "Multiple angles reduce false negatives"]
    bullet_block(s, leg_x, ly, leg_w, bullets,
                 font_size=Pt(12), bullet_color=TEAL)

    return s


def slide_06_system(prs):
    s = add_slide(prs)
    fill_bg(s, WHITE)
    top_bar(s, 6)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "System Architecture")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "What's Inside the System", font_size=Pt(30))

    # Block diagram row
    bd_y = CONTENT_TOP + Inches(1.1)
    bd_h = Inches(2.8)
    box_w = Inches(2.0)
    arrow_w = Inches(0.6)
    hub_w = Inches(1.8)
    pc_w = Inches(2.2)
    out_w = Inches(1.8)

    # --- WET ZONE: 4 camera boxes stacked ---
    wz_label_y = bd_y - Inches(0.28)
    add_textbox(s, ML, wz_label_y, Inches(2.1), Inches(0.25),
                "WET ZONE", font_size=Pt(8), bold=True, color=TEAL_DARK)
    cam_h = Inches(0.55)
    cam_gap = Inches(0.08)
    cam_names = ["Camera A", "Camera B", "Camera C", "Camera D"]
    for i, cam in enumerate(cam_names):
        cy = bd_y + i * (cam_h + cam_gap)
        cb = add_rect(s, ML, cy, box_w, cam_h,
                      fill_color=RGBColor(220, 242, 240),
                      line_color=TEAL, line_width=Pt(1))
        add_textbox(s, ML + Inches(0.1), cy + Inches(0.12),
                    box_w - Inches(0.2), cam_h - Inches(0.2),
                    cam, font_size=Pt(12), bold=False, color=TEAL_DARK)

    # Arrow: cameras -> hub
    ax1 = ML + box_w
    add_textbox(s, ax1, bd_y + Inches(1.0), arrow_w, Inches(0.6),
                "→\nUSB", font_size=Pt(12), color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

    # --- DRY BAY ---
    dry_x = ax1 + arrow_w
    add_textbox(s, dry_x, wz_label_y, Inches(4.2), Inches(0.25),
                "DRY BAY (SEALED)", font_size=Pt(8), bold=True, color=TEAL_DARK)

    hub_y = bd_y + Inches(0.5)
    hub_box = add_rect(s, dry_x, hub_y, hub_w, Inches(0.8),
                       fill_color=RGBColor(255, 246, 220),
                       line_color=GOLDEN, line_width=Pt(1.5))
    add_textbox(s, dry_x + Inches(0.1), hub_y + Inches(0.18),
                hub_w - Inches(0.2), Inches(0.5),
                "Powered USB Hub", font_size=Pt(12), bold=True, color=RGBColor(138, 96, 0))

    pc_y = hub_y + Inches(1.0)
    pc_box = add_rect(s, dry_x, pc_y, pc_w, Inches(1.0),
                      fill_color=TEAL)
    pc_box.line.fill.background()
    add_textbox(s, dry_x + Inches(0.1), pc_y + Inches(0.2),
                pc_w - Inches(0.2), Inches(0.6),
                "Mini PC\n(edge compute)", font_size=Pt(12), bold=True, color=WHITE)

    pwr_y = pc_y + Inches(1.2)
    add_rect(s, dry_x, pwr_y, pc_w, Inches(0.65),
             fill_color=FOG, line_color=SLATE_LIGHT, line_width=Pt(1))
    add_textbox(s, dry_x + Inches(0.1), pwr_y + Inches(0.15),
                pc_w - Inches(0.2), Inches(0.4),
                "Power distribution", font_size=Pt(12), color=SLATE)

    # Arrow: dry bay -> outputs
    ax2 = dry_x + pc_w
    add_textbox(s, ax2, bd_y + Inches(1.0), arrow_w, Inches(0.6),
                "→\nEthernet", font_size=Pt(11), color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

    # --- OUTPUTS ---
    out_x = ax2 + arrow_w
    add_textbox(s, out_x, wz_label_y, Inches(3), Inches(0.25),
                "BUILDING INFRASTRUCTURE", font_size=Pt(8), bold=True, color=SLATE_LIGHT)

    server_box = add_rect(s, out_x, bd_y + Inches(0.4), out_w, Inches(0.9),
                          fill_color=SLATE)
    server_box.line.fill.background()
    add_textbox(s, out_x + Inches(0.1), bd_y + Inches(0.55),
                out_w - Inches(0.2), Inches(0.6),
                "USEFULL\nServer", font_size=Pt(12), bold=True, color=WHITE)

    pwr2_y = bd_y + Inches(1.6)
    add_rect(s, out_x, pwr2_y, out_w, Inches(0.7),
             fill_color=FOG, line_color=SLATE_LIGHT, line_width=Pt(1))
    add_textbox(s, out_x + Inches(0.1), pwr2_y + Inches(0.18),
                out_w - Inches(0.2), Inches(0.4),
                "Facility power\n(power in)", font_size=Pt(11), color=SLATE)

    # Two callouts at the bottom
    cb_y = bd_y + bd_h + Inches(0.15)
    col_w2 = (CONTENT_W - Inches(0.3)) / 2

    add_teal_accent_bar(s, ML, cb_y, Inches(0.7))
    add_body(s, ML + Inches(0.15), cb_y, col_w2 - Inches(0.15), Inches(0.7),
             "Wet zone: only optics. Cameras behind polycarbonate windows — no circuitry exposed to moisture.",
             font_size=Pt(12), color=SLATE)

    rx = ML + col_w2 + Inches(0.3)
    add_teal_accent_bar(s, rx, cb_y, Inches(0.7))
    add_body(s, rx + Inches(0.15), cb_y, col_w2 - Inches(0.15), Inches(0.7),
             "Field-serviceable: swap compute node or USB hub without touching camera boards. Single Ethernet out, single power in.",
             font_size=Pt(12), color=SLATE)

    return s


def slide_07_modularity(prs):
    s = add_slide(prs)
    fill_bg(s, FOG)
    top_bar(s, 7)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Design Constraint")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "The Modularity Challenge", font_size=Pt(30))

    # Left: image
    img_x = ML
    img_y = CONTENT_TOP + Inches(1.05)
    img_w = Inches(5.5)
    img_h = Inches(4.8)

    add_textbox(s, img_x, img_y - Inches(0.28), img_w, Inches(0.25),
                "USEFULL RETURN STATION CONCEPT",
                font_size=Pt(8), bold=True, color=TEAL)

    if os.path.exists(STATION_IMG):
        s.shapes.add_picture(STATION_IMG, img_x, img_y, width=img_w, height=img_h)
    else:
        ph = add_rect(s, img_x, img_y, img_w, img_h,
                      fill_color=RGBColor(230, 230, 230),
                      line_color=SLATE_LIGHT, line_width=Pt(1))
        add_textbox(s, img_x, img_y + Inches(2.0), img_w, Inches(0.5),
                    "[return-station-concept.png]",
                    font_size=Pt(11), color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

    add_textbox(s, img_x, img_y + img_h + Inches(0.05), img_w, Inches(0.3),
                "14\" slot, controlled geometry, designed for this system",
                font_size=Pt(10), color=SLATE_LIGHT, italic=True)

    # Right: bullets + callout
    rx = img_x + img_w + Inches(0.4)
    rw = W - rx - MR
    ry = img_y

    bullet_block(s, rx, ry, rw,
                 ["This concept was designed for USEFULL's own return station — a controlled, known geometry.",
                  "NAU may want to use their existing bins. Different geometries, different slot sizes, different internals.",
                  "A fully modular \"drop-in\" system that works in any bin is significantly harder and more expensive."],
                 font_size=Pt(13), row_height=Inches(0.7))

    reco_y = ry + Inches(2.4)
    reco_box = add_rect(s, rx, reco_y, rw, Inches(1.15),
                        fill_color=RGBColor(255, 246, 220),
                        line_color=GOLDEN, line_width=Pt(1.5))
    add_textbox(s, rx + Inches(0.15), reco_y + Inches(0.12), rw - Inches(0.3), Inches(0.3),
                "RECOMMENDATION", font_size=Pt(9), bold=True, color=RGBColor(138, 96, 0))
    add_textbox(s, rx + Inches(0.15), reco_y + Inches(0.42), rw - Inches(0.3), Inches(0.6),
                "Target one known bin design to prove the concept before attempting modularity.",
                font_size=Pt(13), color=SLATE)

    # Gap visual
    gap_y = reco_y + Inches(1.35)
    hw = (rw - Inches(0.5)) / 2
    add_rect(s, rx, gap_y, hw, Inches(0.5),
             fill_color=RGBColor(220, 242, 240), line_color=TEAL, line_width=Pt(1))
    add_textbox(s, rx, gap_y + Inches(0.1), hw, Inches(0.35),
                "Known geometry", font_size=Pt(11), bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    add_textbox(s, rx + hw, gap_y, Inches(0.5), Inches(0.5),
                "↔", font_size=Pt(16), color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

    add_rect(s, rx + hw + Inches(0.5), gap_y, hw, Inches(0.5),
             fill_color=FOG, line_color=SLATE_LIGHT, line_width=Pt(1))
    add_textbox(s, rx + hw + Inches(0.5), gap_y + Inches(0.1), hw, Inches(0.35),
                "Unknown bins", font_size=Pt(11), bold=True,
                color=SLATE, align=PP_ALIGN.CENTER)

    add_textbox(s, rx, gap_y + Inches(0.55), rw, Inches(0.25),
                "modularity gap", font_size=Pt(10), color=SLATE_LIGHT,
                italic=True, align=PP_ALIGN.CENTER)

    return s


def slide_08_cost(prs):
    s = add_slide(prs)
    fill_bg(s, WHITE)
    top_bar(s, 8)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Budget Reality")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "What Would This Cost?", font_size=Pt(30))

    # Callout
    co_y = CONTENT_TOP + Inches(1.05)
    add_golden_accent_bar(s, ML, co_y, Inches(0.7))
    add_body(s, ML + Inches(0.15), co_y, CONTENT_W - Inches(0.15), Inches(0.7),
             "Joshua's initial estimate: $10k–$15k all-in for his work (4–5 months, ~60 days on-site). "
             "We're planning around ~$8k for his labor. This does NOT include hardware, fabrication, "
             "USEFULL software integration, or the dev workstation.",
             font_size=Pt(12), color=SLATE, italic=True)

    # Table
    from pptx.util import Inches as I
    tbl_y = co_y + Inches(0.85)
    tbl_h = Inches(3.5)

    rows_data = [
        ("Hardware + installation",                    "$3.5k – $6.1k",  "$15.7k – $27.4k"),
        ("Housing / rail design + fabrication",        "$6k – $15k",     "$6k – $15k"),
        ("Dev workstation",                            "$2.5k – $4k",    "$2.5k – $4k"),
        ("USEFULL software integration (Nadia + Yulia)","$15k – $27k",  "$15k – $27k"),
        ("Consultant labor (Joshua)",                  "~$8k",           "~$8k"),
        ("All-in with 20% contingency",                "~$36k – $65k",  "~$50k – $89k"),
    ]

    col_widths = [Inches(7.2), Inches(2.2), Inches(2.6)]
    tbl = s.shapes.add_table(len(rows_data) + 1, 3,
                              ML, tbl_y, sum(col_widths), tbl_h).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    # Header row
    headers = ["Category", "2-Bin Pilot", "9-Bin Full Rollout"]
    for ci, h_text in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = h_text
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT

    # Data rows
    for ri, (cat, pilot, full) in enumerate(rows_data, start=1):
        is_total = ri == len(rows_data)
        row_fill = RGBColor(220, 242, 240) if is_total else (
            RGBColor(250, 250, 250) if ri % 2 == 0 else WHITE
        )
        for ci, val in enumerate([cat, pilot, full]):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(12)
            run.font.bold = is_total
            run.font.color.rgb = TEAL_DARK if is_total else SLATE
            run.font.name = FONT

    # Note
    note_y = tbl_y + tbl_h + Inches(0.15)
    add_body(s, ML, note_y, CONTENT_W, Inches(0.4),
             "Camera electronics (~$400/bin) are not the expensive part. Software integration, packaging, and iteration are. "
             "Recurring cloud/inference: est. $500–$1,200/month if hosted inference is used later.",
             font_size=Pt(10), color=SLATE_LIGHT, italic=True)

    return s


def slide_09_integration(prs):
    s = add_slide(prs)
    fill_bg(s, FOAM)
    top_bar(s, 9)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "The Hidden Bulk")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "USEFULL Software Integration — The Largest Cost", font_size=Pt(28))

    add_body(s, ML, CONTENT_TOP + Inches(1.0), CONTENT_W, Inches(0.35),
             "Even if the camera system works perfectly, significant dev effort is still required to make it operationally useful.",
             font_size=Pt(13), color=SLATE_LIGHT, italic=True)

    # Left: hours table
    tbl_y = CONTENT_TOP + Inches(1.5)
    col_widths = [Inches(5.8), Inches(1.2)]
    rows_data = [
        ("Integration design + interface definition",   "10–18"),
        ("Event ingestion pipeline",                    "14–25"),
        ("DB / schema work",                            "8–17"),
        ("Scan ↔ return matching logic (hardest)",      "28–49"),
        ("\"Scanned but not dropped\" workflows",       "18–35"),
        ("Ops / support tooling",                       "14–28"),
        ("Logging / observability",                     "14–25"),
        ("QA support + bug fixing + rollout",           "21–39"),
        ("Knowledge transfer / coordination",           "10–18"),
        ("Total Nadia",                                 "137–254"),
    ]

    tbl_h = Inches(4.2)
    tbl = s.shapes.add_table(len(rows_data) + 1, 2,
                              ML, tbl_y, sum(col_widths), tbl_h).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h_text in enumerate(["Work Package (Nadia)", "Hours"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = h_text
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = WHITE; run.font.name = FONT

    for ri, (pkg, hrs) in enumerate(rows_data, start=1):
        is_total = ri == len(rows_data)
        rf = RGBColor(220, 242, 240) if is_total else (RGBColor(248,252,251) if ri%2==0 else WHITE)
        for ci, val in enumerate([pkg, hrs]):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rf
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(11); run.font.bold = is_total
            run.font.color.rgb = TEAL_DARK if is_total else SLATE
            run.font.name = FONT

    add_textbox(s, ML, tbl_y + tbl_h + Inches(0.08), Inches(7), Inches(0.25),
                "Yulia (QA): 42–84 hours in parallel",
                font_size=Pt(10), color=SLATE_LIGHT, italic=True)

    # Right: cost cards
    card_x = ML + sum(col_widths) + Inches(0.4)
    card_w = W - card_x - MR
    cy = tbl_y

    for (label, amount, detail, top_color) in [
        ("Nadia (Dev)",  "$12k – $23k", "137–254 hrs @ ~$90/hr loaded", TEAL),
        ("Yulia (QA)",   "$1.9k – $3.8k", "42–84 hrs @ ~$45/hr loaded", GOLDEN),
    ]:
        card = add_rect(s, card_x, cy, card_w, Inches(1.2),
                        fill_color=WHITE, line_color=RGBColor(220,220,220), line_width=Pt(1))
        top = add_rect(s, card_x, cy, card_w, Inches(0.06), fill_color=top_color)
        top.line.fill.background()
        add_textbox(s, card_x + Inches(0.12), cy + Inches(0.1), card_w - Inches(0.24), Inches(0.28),
                    label, font_size=Pt(10), bold=True, color=SLATE_LIGHT)
        add_textbox(s, card_x + Inches(0.12), cy + Inches(0.38), card_w - Inches(0.24), Inches(0.45),
                    amount, font_size=Pt(20), bold=True, color=top_color)
        add_textbox(s, card_x + Inches(0.12), cy + Inches(0.85), card_w - Inches(0.24), Inches(0.28),
                    detail, font_size=Pt(10), color=SLATE_LIGHT)
        cy += Inches(1.35)

    # Combined total card
    total_card = add_rect(s, card_x, cy, card_w, Inches(1.5), fill_color=TEAL)
    total_card.line.fill.background()
    add_textbox(s, card_x + Inches(0.15), cy + Inches(0.12), card_w - Inches(0.3), Inches(0.28),
                "COMBINED", font_size=Pt(10), bold=True, color=RGBColor(180,225,227))
    add_textbox(s, card_x + Inches(0.15), cy + Inches(0.42), card_w - Inches(0.3), Inches(0.55),
                "$15k – $27k", font_size=Pt(24), bold=True, color=WHITE)
    add_textbox(s, card_x + Inches(0.15), cy + Inches(1.0), card_w - Inches(0.3), Inches(0.35),
                "Likely range: $16k–$22k", font_size=Pt(11), color=RGBColor(200,235,237))

    return s


def slide_10_options(prs):
    s = add_slide(prs)
    fill_bg(s, WHITE)
    top_bar(s, 10)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Decision Framework")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "Three Options", font_size=Pt(30))

    # 3-column cards
    card_y = CONTENT_TOP + Inches(1.0)
    card_h = Inches(5.6)
    card_w = (CONTENT_W - Inches(0.4)) / 3
    gap = Inches(0.2)

    options = [
        {
            "label": "OPTION A",
            "title": "Engage Consultant — Stage-Gated Pilot",
            "meta": "2-bin pilot at UMass Lowell · ~$36k–$65k all-in · 4–5 months",
            "body": "Start with known bin geometry. Decision gate after pilot: does it work well enough? Only expand to 9 bins if pilot succeeds.",
            "pro": "Most thorough path; real CV expertise; proves concept before scaling",
            "con": "Significant investment before we confirm the problem it solves is driving losses",
            "top": TEAL,
        },
        {
            "label": "OPTION B",
            "title": "Small Internal Pilot — No Consultant",
            "meta": "Off-the-shelf cameras · basic motion detection · lower upfront cost",
            "body": "Use existing team knowledge to prototype a simpler version. Skip custom throat design to start.",
            "pro": "Cheaper to start; keeps optionality; builds internal knowledge",
            "con": "No CV expertise; risk of building something that doesn't validate the concept",
            "top": GOLDEN,
        },
        {
            "label": "OPTION C",
            "title": "Shelve for Now",
            "meta": "Gather loss data first · revisit when loss driver is confirmed",
            "body": "Focus on root causes suggested by NAU and Lowell feedback. Confirm return-station fraud is a significant driver before investing.",
            "pro": "No investment risk; addresses \"are we solving the right problem?\" first",
            "con": "If CEO's instinct is right, we lose time; consultant may not be available later",
            "top": SLATE_LIGHT,
        },
    ]

    for i, opt in enumerate(options):
        cx = ML + i * (card_w + gap)
        card = add_rect(s, cx, card_y, card_w, card_h,
                        fill_color=WHITE,
                        line_color=RGBColor(220, 220, 220), line_width=Pt(1))
        top = add_rect(s, cx, card_y, card_w, Inches(0.07), fill_color=opt["top"])
        top.line.fill.background()

        add_textbox(s, cx + Inches(0.15), card_y + Inches(0.12), card_w - Inches(0.3), Inches(0.28),
                    opt["label"], font_size=Pt(10), bold=True, color=opt["top"])
        add_textbox(s, cx + Inches(0.15), card_y + Inches(0.42), card_w - Inches(0.3), Inches(0.65),
                    opt["title"], font_size=Pt(14), bold=True, color=SLATE)
        add_textbox(s, cx + Inches(0.15), card_y + Inches(1.1), card_w - Inches(0.3), Inches(0.5),
                    opt["meta"], font_size=Pt(10), color=SLATE_LIGHT, italic=True)
        add_textbox(s, cx + Inches(0.15), card_y + Inches(1.65), card_w - Inches(0.3), Inches(1.2),
                    opt["body"], font_size=Pt(12), color=SLATE)

        # Pro/Con
        pro_y = card_y + Inches(3.0)
        add_teal_accent_bar(s, cx + Inches(0.15), pro_y, Inches(0.5), w=Inches(0.04))
        add_textbox(s, cx + Inches(0.28), pro_y, card_w - Inches(0.4), Inches(0.55),
                    "+" + " " + opt["pro"], font_size=Pt(11), color=TEAL_DARK)

        con_y = pro_y + Inches(0.65)
        add_golden_accent_bar(s, cx + Inches(0.15), con_y, Inches(0.5), w=Inches(0.04))
        add_textbox(s, cx + Inches(0.28), con_y, card_w - Inches(0.4), Inches(0.55),
                    "–" + " " + opt["con"], font_size=Pt(11), color=SLATE_LIGHT)

    return s


def slide_11_unknowns(prs):
    s = add_slide(prs)
    fill_bg(s, FOAM)
    top_bar(s, 11)

    add_label(s, ML, CONTENT_TOP, CONTENT_W, "Before We Decide")
    add_heading(s, ML, CONTENT_TOP + Inches(0.32), CONTENT_W,
                "Key Unknowns & Recommended Next Steps", font_size=Pt(28))

    col_w = (CONTENT_W - Inches(0.4)) / 2
    col_y = CONTENT_TOP + Inches(1.05)

    # Left: unknowns
    add_textbox(s, ML, col_y, col_w, Inches(0.3),
                "WHAT WE DON'T KNOW YET",
                font_size=Pt(9), bold=True, color=SLATE_LIGHT)

    unknowns = [
        "Whether return-station fraud is actually a significant source of container loss (NAU + Lowell evidence suggests maybe not)",
        "Whether the no-illumination approach works in real dining hall lighting on wet, reflective surfaces",
        "Whether NAU's existing bins can accommodate the system — or if they'd adopt USEFULL return stations",
        "Whether Nadia's bandwidth can absorb this without impacting the roadmap",
        "Exact mechanical fit — 1.75\" rails are tight once fasteners, brackets, and cable bends are added",
    ]

    uy = col_y + Inches(0.4)
    for u in unknowns:
        # golden question square
        q = add_rect(s, ML, uy + Inches(0.05), Inches(0.18), Inches(0.18),
                     fill_color=RGBColor(255, 246, 220),
                     line_color=GOLDEN, line_width=Pt(1))
        add_textbox(s, ML + Inches(0.22), uy, col_w - Inches(0.25), Inches(0.55),
                    u, font_size=Pt(12), color=SLATE)
        uy += Inches(0.62)

    # Right: next steps + bottom line
    rx = ML + col_w + Inches(0.4)
    rw = W - rx - MR

    add_textbox(s, rx, col_y, rw, Inches(0.3),
                "REGARDLESS OF OPTION — IMMEDIATE NEXT STEPS",
                font_size=Pt(9), bold=True, color=SLATE_LIGHT)

    steps = [
        "Quantify the actual loss rate at return stations specifically, vs. other points in the container lifecycle",
        "Clarify whether NAU needs to use existing bins or would adopt USEFULL return stations",
        "If proceeding: start with a 2-bin pilot at one site with known geometry — don't try to solve modularity first",
    ]

    sy = col_y + Inches(0.4)
    for i, step in enumerate(steps, 1):
        num_box = add_rect(s, rx, sy + Inches(0.04), Inches(0.28), Inches(0.28),
                           fill_color=TEAL)
        num_box.line.fill.background()
        add_textbox(s, rx, sy + Inches(0.04), Inches(0.28), Inches(0.28),
                    str(i), font_size=Pt(11), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, rx + Inches(0.38), sy, rw - Inches(0.38), Inches(0.7),
                    step, font_size=Pt(12), color=SLATE)
        sy += Inches(0.8)

    # Bottom line callout
    bl_y = sy + Inches(0.2)
    bl_box = add_rect(s, rx, bl_y, rw, Inches(1.35),
                      fill_color=RGBColor(220, 242, 240),
                      line_color=TEAL, line_width=Pt(1.5))
    add_teal_accent_bar(s, rx, bl_y, Inches(1.35), w=Inches(0.06))
    add_body(s, rx + Inches(0.2), bl_y + Inches(0.18), rw - Inches(0.28), Inches(1.0),
             "A camera-based return verification system is technically plausible. The question is whether "
             "it solves the right problem, and whether the investment is justified given what we're learning "
             "about where containers are actually being lost.",
             font_size=Pt(12), color=TEAL_DARK, italic=True)

    return s


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def build():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_right_problem(prs)
    slide_04_joshua(prs)
    slide_05_throat(prs)
    slide_06_system(prs)
    slide_07_modularity(prs)
    slide_08_cost(prs)
    slide_09_integration(prs)
    slide_10_options(prs)
    slide_11_unknowns(prs)

    out = "return-verification.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
